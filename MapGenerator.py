import folium
import pandas as pd
import geopandas
import webbrowser
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from selenium import webdriver
from PIL import Image
import time
import os
from pptx import Presentation
from pptx.util import Inches

# Function to read site info from a file
def readSiteInfo(filePath):
    df = pd.read_csv(filePath)
    required_columns = ["SiteName", "Lat", "Long"]

    if not all(column in df.columns for column in required_columns):
        missing_columns = [column for column in required_columns if column not in df.columns]
        raise ValueError(f"Missing columns in site info file: {', '.join(missing_columns)}")

    df1 = df.loc[:, required_columns]
    geoSite = geopandas.points_from_xy(df1.Long, df1.Lat)
    geoSiteDF = geopandas.GeoDataFrame(df1[required_columns], geometry=geoSite)
    geoSiteList = [[point.xy[1][0], point.xy[0][0]] for point in geoSiteDF.geometry]
    return geoSiteList, geoSiteDF

# Function to read path info from a file
def readPathInfo(filePath):
    df = pd.read_csv(filePath)
    required_columns = ["Latitude", "Longitude"]

    if not all(column in df.columns for column in required_columns):
        missing_columns = [column for column in required_columns if column not in df.columns]
        raise ValueError(f"Missing columns in path info file: {', '.join(missing_columns)}")

    df1 = df.loc[:, required_columns]
    geoPath = geopandas.points_from_xy(df1.Longitude, df1.Latitude)
    geoPathDF = geopandas.GeoDataFrame(df1[required_columns], geometry=geoPath)
    geoPathList = [[point.xy[1][0], point.xy[0][0]] for point in geoPathDF.geometry]
    return geoPathList, geoPathDF

# Function to create a folium map
def createMap(geoSiteList):
    map_center = geoSiteList[0] if geoSiteList else [0, 0]
    return folium.Map(location=map_center, zoom_start=12)

# Function to add site markers to the map
def addSiteMarkers(geoSiteList, geoSiteDF, map, site_color='red'):
    for i, location in enumerate(geoSiteList):
        if site_color in folium.Icon.color_options:
            folium.Marker(location=location, popup="SiteName: " + str(geoSiteDF.SiteName[i]),
                          icon=folium.Icon(color=site_color)).add_to(map)
        else:
            folium.Marker(location=location, popup="SiteName: " + str(geoSiteDF.SiteName[i]),
                          icon=folium.DivIcon(html=f"""
                              <div style="background-color:{site_color}; width:12px; height:12px;
                              border-radius:50%; border:1px solid black;"></div>""")).add_to(map)

def addPathMarkers(geoPathList, map, path_color='black'):
    for location in geoPathList:
        folium.CircleMarker(location=location, radius=1, color=path_color, fill=True, fill_color=path_color).add_to(map)


# Function to fit map bounds
def fitMapBounds(map, geoSiteList, geoPathList):
    all_locations = geoSiteList + geoPathList
    map.fit_bounds(all_locations)

# Function to save and open the map in a web browser
def openMap(map):
    map.save("map.html")
    webbrowser.open("map.html")

# Function to save the map as a PNG
def saveMapAsImage(map):
    map.save("map.html")

    # Set up selenium WebDriver (you may need to adjust the path to your chromedriver)
    options = webdriver.ChromeOptions()
    options.add_argument('--headless')
    options.add_argument('--no-sandbox')
    options.add_argument('--disable-dev-shm-usage')
    driver = webdriver.Chrome(options=options)

    # Open the HTML file in the headless browser
    driver.get("file://" + os.path.abspath("map.html"))
    time.sleep(5)  

    # Save the screenshot
    driver.save_screenshot("map.png")
    driver.quit()
    saveImageToPowerPoint("map.png", "map_presentation.pptx")

def saveImageToPowerPoint(image_path, pptx_path):
    # Create a presentation object
    prs = Presentation()

    # Add a slide
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add image to the slide
    left = Inches(0)
    top = Inches(0)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(10), height=Inches(7.5))

    # Save the presentation
    prs.save(pptx_path)

# Tkinter GUI setup
class App:
    def __init__(self, root):
        self.root = root
        self.root.title("Map Generator")
        self.root.configure(bg='#5083a0')

        self.siteFilePath = ""
        self.pathFilePath = ""

        self.create_widgets()

    def create_widgets(self):
        # Load image
        img = tk.PhotoImage(file="logo.png")
        img_label = ttk.Label(self.root, image=img, background='#5083a0')
        img_label.image = img  # Keep a reference to the image to prevent garbage collection
        img_label.pack(pady=10)

        frame = ttk.Frame(self.root, padding="20", style='TFrame')
        frame.pack(padx=10, pady=10, fill='x', expand=True)

        import_site_btn = ttk.Button(frame, text="Select Site Info File", command=self.select_site_file, style='TButton')
        import_site_btn.pack(pady=10, fill='x')

        import_path_btn = ttk.Button(frame, text="Select Path Info File", command=self.select_path_file, style='TButton')
        import_path_btn.pack(pady=10, fill='x')

        self.file_label = ttk.Label(frame, text="No files selected", style='TLabel')
        self.file_label.pack(padx=50, pady=10, fill='x')

        create_map_btn = ttk.Button(frame, text="Create Map", command=self.create_map, style='TButton')
        create_map_btn.pack(pady=10, fill='x')

        save_image_btn = ttk.Button(frame, text="Save Map as Image", command=self.save_map_as_image, style='TButton')
        save_image_btn.pack(pady=10, fill='x')

    def select_site_file(self):
        self.siteFilePath = filedialog.askopenfilename(filetypes=[("CSV files", "*.csv")])
        if self.siteFilePath:
            self.update_file_label()
            messagebox.showinfo("Selected File", f"Site info file selected: {self.siteFilePath}")

    def select_path_file(self):
        self.pathFilePath = filedialog.askopenfilename(filetypes=[("CSV files", "*.csv")])
        if self.pathFilePath:
            self.update_file_label()
            messagebox.showinfo("Selected File", f"Path info file selected: {self.pathFilePath}")

    def update_file_label(self):
        site_file = self.siteFilePath.split('/')[-1] if self.siteFilePath else "None"
        path_file = self.pathFilePath.split('/')[-1] if self.pathFilePath else "None"
        self.file_label.config(text=f"Site File: {site_file}\n Path File: {path_file}")

    def create_map(self):
        if not self.siteFilePath or not self.pathFilePath:
            messagebox.showwarning("Warning", "Please select both site info and path info files")
            return

        try:
            geoSiteList, geoSiteDF = readSiteInfo(self.siteFilePath)
            geoPathList, geoPathDF = readPathInfo(self.pathFilePath)

            map = createMap(geoSiteList)
            addSiteMarkers(geoSiteList, geoSiteDF, map)
            addPathMarkers(geoPathList, map)
            fitMapBounds(map, geoSiteList, geoPathList)
            openMap(map)
        except ValueError as e:
            messagebox.showerror("Error", str(e))

    def save_map_as_image(self):
        if not self.siteFilePath or not self.pathFilePath:
            messagebox.showwarning("Warning", "Please select both site info and path info files")
            return

        try:
            geoSiteList, geoSiteDF = readSiteInfo(self.siteFilePath)
            geoPathList, geoPathDF = readPathInfo(self.pathFilePath)

            map = createMap(geoSiteList)
            addSiteMarkers(geoSiteList, geoSiteDF, map)
            addPathMarkers(geoPathList, map)
            fitMapBounds(map, geoSiteList, geoPathList)
            saveMapAsImage(map)
            messagebox.showinfo("Save Complete", "Map has been saved as map.png")
        except ValueError as e:
            messagebox.showerror("Error", str(e))
